#!/usr/bin/env python3
"""
Convert RAG Benchmark content to PowerPoint (.pptx) - Enhanced with detailed RAG intro
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "RAG Benchmark Suite",
        "Quantifying RAG Performance on Documents & Code\n\n"
        "Part 1: RAG Fundamentals (Slides 2-6)\n"
        "Part 2: Document Q&A Benchmark (Slides 7-11)\n"
        "Part 3: Code Modification Benchmark (Slides 12-17)\n\n"
        "Kaito Project Team | November 2025")
    
    # Slide 2: What is RAG? - Detailed
    add_content_slide(prs,
        "What is RAG?",
        "Retrieval-Augmented Generation - A Hybrid Approach",
        [
            "The Core Problem with Pure LLMs:",
            "  • Knowledge cutoff date (e.g., training data ends in 2023)",
            "  • Cannot access private/proprietary documents",
            "  • No real-time information (stock prices, news, etc.)",
            "  • Hallucinate facts when uncertain",
            "",
            "What is RAG?",
            "  RAG = Retrieval (search documents) + Augmented (add context)",
            "        + Generation (LLM generates answer)",
            "",
            "  Instead of asking LLM to answer from memory alone,",
            "  RAG first retrieves relevant information, then asks LLM",
            "  to answer based on that retrieved context.",
            "",
            "Simple Analogy:",
            "  Pure LLM = Closed-book exam (rely on memory)",
            "  RAG = Open-book exam (can reference materials)"
        ])
    
    # Slide 3: RAG Components
    add_content_slide(prs,
        "RAG System Components",
        "Four Key Components Working Together",
        [
            "1. Document Loader",
            "   • Ingests various formats (PDF, TXT, DOCX, HTML, Code)",
            "   • Extracts text and metadata",
            "",
            "2. Text Chunker (Splitter)",
            "   • Breaks documents into smaller chunks (512-1024 tokens)",
            "   • Maintains context with overlapping windows (50-100 tokens)",
            "   • Preserves semantic boundaries (paragraphs, sentences)",
            "",
            "3. Embedding Model + Vector Database",
            "   • Embedding Model: Converts text → dense vectors (768D)",
            "     Common: OpenAI ada-002, sentence-transformers",
            "   • Vector Database: Stores embeddings for fast retrieval",
            "     Examples: Faiss, Pinecone, Chroma, Weaviate",
            "   • Enables semantic search (meaning-based)",
            "",
            "4. Retriever + Generator (LLM)",
            "   • Retriever: Finds top-k most relevant chunks",
            "   • Generator: LLM creates answer using retrieved context"
        ])
    
    # Slide 4: RAG Workflow
    add_content_slide(prs,
        "RAG Workflow: From Query to Answer",
        "Step-by-Step Process",
        [
            "Indexing Phase (One-time Setup):",
            "  1. Load documents → 2. Chunk text → 3. Generate embeddings",
            "  4. Store in vector database",
            "",
            "Query Phase (Every User Request):",
            "",
            "  User Query: \"What is the API timeout limit?\"",
            "     ↓",
            "  Step 1: Convert query to embedding vector",
            "     ↓",
            "  Step 2: Search vector database (cosine similarity)",
            "     Retrieve top-5 chunks: [0.92, 0.88, 0.85, ...]",
            "     ↓",
            "  Step 3: Build augmented prompt",
            "     Context: [Retrieved chunks]",
            "     Question: \"What is the API timeout limit?\"",
            "     ↓",
            "  Step 4: LLM generates answer (based on context)",
            "     ↓",
            "  Answer: \"The API timeout is 30 seconds",
            "          according to api_config.yaml\""
        ])
    
    # Slide 5: RAG vs Fine-Tuning
    add_content_slide(prs,
        "RAG vs Fine-Tuning",
        "Two Different Approaches to Customizing LLMs",
        [
            "Fine-Tuning:",
            "  • Retrains model weights on your specific data",
            "  • Teaches model new patterns, style, or domain knowledge",
            "  • Model internalizes knowledge (stored in parameters)",
            "  • One-time training process (expensive, time-consuming)",
            "",
            "RAG (Retrieval-Augmented Generation):",
            "  • Keeps model frozen, adds external knowledge retrieval",
            "  • Provides relevant context at inference time",
            "  • Knowledge stored externally (in vector database)",
            "  • Easy to update (just re-index new documents)",
            "",
            "Key Difference:",
            "  Fine-Tuning = Teaching the model new knowledge",
            "  RAG = Giving the model reference materials to consult"
        ])
    
    # Slide 6: RAG vs Fine-Tuning Table
    add_table_slide(prs,
        "RAG vs Fine-Tuning: Use Cases",
        "Choosing the Right Approach",
        ["Aspect", "RAG", "Fine-Tuning"],
        [
            ["Knowledge Updates", "Easy (re-index)", "Hard (retrain)"],
            ["Data Requirements", "Any amount", "1000s of examples"],
            ["Cost", "Low (inference)", "High (training + GPU)"],
            ["Speed to Deploy", "Hours", "Days to weeks"],
            ["Use Case", "Q&A, search", "Style, reasoning"],
            ["Explainability", "High (sources)", "Low (black box)"],
            ["Accuracy on Facts", "High (grounded)", "Medium"]
        ],
        [
            "",
            "Use RAG when:",
            "  ✓ Need frequently-updated knowledge bases",
            "  ✓ Want to cite sources and transparency",
            "  ✓ Have limited budget and time",
            "  Examples: Customer support, documentation search",
            "",
            "Use Fine-Tuning when:",
            "  ✓ Need to change model behavior or style",
            "  ✓ Want domain-specific reasoning patterns",
            "  ✓ Have sufficient training data and compute",
            "  Examples: Code generation, medical diagnosis"
        ])
    
    # Slide 7: Why Benchmark?
    add_content_slide(prs,
        "Why We Need RAG Benchmarks",
        "\"How much better is RAG compared to pure LLM?\"",
        [
            "Without Benchmarks:",
            "  ❓ Unclear if RAG adds value",
            "  ❓ Don't know optimal configuration",
            "  ❓ Hard to justify investment",
            "",
            "With Benchmarks:",
            "  ✓ Quantitative metrics: Success rate, accuracy scores",
            "  ✓ Cost analysis: Token usage, API costs",
            "  ✓ A/B comparison: RAG vs Baseline side-by-side",
            "  ✓ Data-driven decisions: Prove ROI with numbers",
            "",
            "Our Solution: Two Specialized Benchmarks",
            "  1. Document Q&A: For documents, PDFs, manuals",
            "  2. Code Modification: For bug fixes, features"
        ])
    
    # Slide 8: Document Benchmark Overview
    add_content_slide(prs,
        "RAG Benchmark for Documents",
        "Measure RAG performance on document-based Q&A",
        [
            "What It Tests:",
            "  📚 Document retrieval accuracy",
            "  ✅ Answer quality: RAG vs pure LLM (factual)",
            "  🧠 Comprehension: RAG vs pure LLM (analytical)",
            "  💰 Token efficiency: Cost comparison",
            "",
            "Key Features:",
            "  • Generates 20 test questions automatically",
            "  • Tests both RAG and pure LLM on same questions",
            "  • Uses LLM-as-Judge for scoring (0-10 scale)",
            "  • Produces detailed reports with metrics",
            "",
            "Typical Results:",
            "  RAG Average Score:  8.5/10  (+89% improvement)",
            "  Pure LLM Score:     4.5/10",
            "  Token Usage:        -15% (RAG more efficient)"
        ])
    
    # Slide 9: Document Workflow
    add_content_slide(prs,
        "Document Benchmark Workflow",
        "5-Step Process",
        [
            "PREREQUISITE: User indexes documents in RAG system",
            "",
            "STEP 1: Generate Test Questions",
            "  • Query RAG index to retrieve 20 content nodes",
            "  • LLM generates Q&A pairs from each node",
            "  • 10 closed (factual) + 10 open (analytical)",
            "",
            "STEP 2: Run RAG System",
            "  • For each question: search → retrieve context",
            "  • LLM generates answer using context",
            "",
            "STEP 3: Run Pure LLM (No RAG)",
            "  • Same questions, no document access",
            "  • LLM relies on pre-trained knowledge only",
            "",
            "STEP 4: LLM Judge Evaluation",
            "  • Judge LLM scores each answer (0-10)",
            "",
            "STEP 5: Generate Comparison Report",
            "  • Average scores, improvement percentage"
        ])
    
    # Slide 10: Question Types
    add_content_slide(prs,
        "Two Question Types",
        "Comprehensive Testing",
        [
            "Closed Questions (Factual Accuracy)",
            "  Definition: Specific, verifiable answers",
            "  Examples:",
            "    - \"What is the maximum timeout for API requests?\"",
            "    - \"Which port does the service listen on?\"",
            "  Scoring (0/5/10):",
            "    10: Completely correct",
            "    5: Partially correct",
            "    0: Wrong or irrelevant",
            "",
            "Open Questions (Comprehension & Analysis)",
            "  Definition: Understanding and synthesis required",
            "  Examples:",
            "    - \"How does the system handle concurrent requests?\"",
            "    - \"Explain the error handling strategy\"",
            "  Scoring (0-10 gradient):",
            "    Accuracy (3) + Completeness (3) +",
            "    Understanding (2) + Relevance (2)"
        ])
    
    # Slide 11: Document Results
    add_table_slide(prs,
        "Document Benchmark - Results",
        "Real Performance Comparison",
        ["Metric", "RAG System", "Pure LLM", "Improvement"],
        [
            ["Overall Score", "8.5/10", "4.5/10", "+89%"],
            ["Closed Questions", "9.2/10", "3.8/10", "+142%"],
            ["Open Questions", "7.8/10", "5.2/10", "+50%"],
            ["Token Usage", "45K", "53K", "-15%"]
        ],
        [
            "",
            "Key Findings:",
            "  ✓ RAG excels at factual questions (+142%)",
            "  ✓ RAG improves comprehension (+50%)",
            "  ✓ RAG is more token-efficient (-15%)",
            "  ❌ Pure LLM struggles without document access"
        ])
    
    # Slide 12: Code Benchmark Overview
    add_content_slide(prs,
        "RAG Benchmark for Code Modification",
        "Measure RAG performance on automated code fixes",
        [
            "What It Tests:",
            "  🐛 Bug fixing accuracy: Success rate",
            "  ✅ Test validation: Validated through unit tests",
            "  💰 Token efficiency: With TOP-4 filtering",
            "  📁 File selection: RAG auto vs manual context",
            "",
            "Key Difference from Document Benchmark:",
            "  • Document: Evaluate answers with LLM judge",
            "  • Code: Validate with unit tests (objective)",
            "",
            "Typical Results:",
            "  Baseline (Manual):  20% success (1/5 issues)",
            "  RAG (Automatic):     0% success (0/5 issues)",
            "  Token Savings:      21.6% (with TOP-4 filtering)",
            "",
            "Note: Benchmark identifies RAG limitations"
        ])
    
    # Slide 13: Code Workflow
    add_content_slide(prs,
        "Code Benchmark Workflow",
        "4-Step Process",
        [
            "PREREQUISITE: Index Code Repository",
            "  python rag.py --repo . --index code_repo_benchmark",
            "",
            "STEP 1: Generate Test Issues",
            "  • Scan repository structure",
            "  • Identify components (packages, modules)",
            "  • Generate realistic issues (5-10)",
            "",
            "STEP 2: Run Baseline Solution (Manual)",
            "  • Developer provides relevant file list",
            "  • LLM modifies code with manual context",
            "  • Apply changes → Run tests → Pass/Fail",
            "",
            "STEP 3: Run RAG Solution (Automatic)",
            "  • RAG retrieves 100+ files internally",
            "  • TOP-4 Filter: Sort by relevance, take top 4",
            "  • Apply changes → Run tests → Pass/Fail",
            "",
            "STEP 4: Compare Results",
            "  • Success rate, token efficiency, error analysis"
        ])
    
    # Slide 14: TOP-4 Filtering
    add_content_slide(prs,
        "TOP-4 Relevance Filtering",
        "Key Innovation",
        [
            "The Problem:",
            "  • RAG retrieves 100+ documents internally",
            "  • Returns 4-16 source_nodes with scores",
            "  • Too many files = token bloat + confusion",
            "",
            "Our Solution: TOP-4 Filtering",
            "  file_scores = {",
            "    \"workspace_validation.go\": 0.5205,",
            "    \"workspace_types.go\": 0.4962,",
            "    \"workspace_controller.go\": 0.4751,",
            "    \"workspace_service.go\": 0.4683,",
            "    \"workspace_test.go\": 0.4512,  # Filtered",
            "  }",
            "  top_4 = sorted_files[:4]",
            "",
            "Results:",
            "  ✓ 21.6% token savings",
            "  ✓ Reduced context confusion",
            "  ✓ Faster LLM processing"
        ])
    
    # Slide 15: Test Validation
    add_content_slide(prs,
        "Objective Validation with Unit Tests",
        "Code uses objective tests (unlike Document's LLM judge)",
        [
            "Validation Process:",
            "",
            "1. Apply Code Modifications",
            "   • Write changed files, backup originals",
            "",
            "2. Generate Git Diff",
            "   • git diff > issue_001.diff",
            "",
            "3. Run Unit Tests",
            "   • go test ./... (Go) or pytest (Python)",
            "   • Capture stdout/stderr",
            "",
            "4. Pass or Fail?",
            "   PASS → Keep changes",
            "   FAIL → Revert changes",
            "",
            "Pass/Fail Criteria:",
            "  ✅ PASS = All tests pass + No compilation errors",
            "  ❌ FAIL = Any test fails OR Compilation error"
        ])
    
    # Slide 16: Code Results
    add_table_slide(prs,
        "Code Benchmark - Results",
        "Real-World Performance on Kaito Repository",
        ["Metric", "Baseline", "RAG", "Notes"],
        [
            ["Success Rate", "20% (1/5)", "0% (0/5)", "RAG needs work"],
            ["Avg Tokens/Issue", "12,543", "9,842", "-21.6% tokens"],
            ["Files Modified", "3-4 files", "4 files", "TOP-4 filter"],
            ["Compilation Errors", "0", "2", "Structure deletion"],
            ["Test Failures", "4", "3", "Logic errors"]
        ],
        [
            "",
            "Key Findings:",
            "  ✓ TOP-4 filtering works perfectly",
            "  ❌ RAG struggles with code structure",
            "  ✓ Manual context wins (for now)",
            "",
            "Action Items:",
            "  1. Strengthen system prompts",
            "  2. Try GPT-4, Claude-3",
            "  3. Improve RAG retrieval quality"
        ])
    
    # Slide 17: Summary
    add_content_slide(prs,
        "Summary & Comparison",
        "Two Benchmarks, Two Use Cases",
        [
            "Document Q&A Benchmark",
            "  Best For: PDFs, reports, documentation, Q&A",
            "  Results:",
            "    ✅ RAG wins (+89% improvement)",
            "    ✅ Strong on factual questions",
            "    ✅ Token efficient",
            "",
            "Code Modification Benchmark",
            "  Best For: Bug fixes, feature additions",
            "  Results:",
            "    ⚠️ Baseline wins (20% vs 0%)",
            "    ✅ TOP-4 filtering saves 21.6% tokens",
            "    ❌ RAG needs improvement",
            "",
            "Key Takeaways:",
            "  1. RAG is powerful for documents (+89%)",
            "  2. RAG needs work for code (0% success)",
            "  3. Benchmarks provide data for decisions",
            "  4. TOP-4 filtering balances quality & efficiency"
        ])
    
    return prs

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(54)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_content_slide(prs, title, subtitle, content):
    """Add content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)
    
    # Content
    top = Inches(1.6) if subtitle else Inches(1.2)
    content_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.space_before = Pt(4)
        if line.startswith("  "):
            p.level = 1

def add_table_slide(prs, title, subtitle, headers, rows, footer):
    """Add table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
    
    # Table
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(1), Inches(1.8), Inches(8), Inches(2.5)
    ).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2.5))
        tf = footer_box.text_frame
        for i, line in enumerate(footer):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)

def main():
    print("🎨 Creating enhanced PowerPoint presentation...")
    prs = create_presentation()
    output = "RAG_Benchmark_Presentation.pptx"
    prs.save(output)
    print(f"✅ Presentation created: {output}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
